import os
import re
import requests
import logging
from io import BytesIO
from typing import Optional, Dict, Any, List
from pptx import Presentation
from pptx.presentation import Presentation as PresentationObject
from pptx.chart.data import CategoryChartData
from app.domain.utils import get_class_name
from app.domain.repository.file_document_repository import IFileDocumentRepository

class PptManager(IFileDocumentRepository):
    def __init__(self, document_path: str, is_remote: Optional[bool] = False):
        self.document_path = document_path
        self.is_remote = is_remote
        self.logger = logging.getLogger(get_class_name(self))
        self.pattern = r"\{([A-Z0-9_]+)\}"
        
    def get_presentation(self) -> Presentation:
        """Obtiene el objeto Presentation desde una URL remota o una ruta local."""
        if self.is_remote:
            response = requests.get(self.document_path)
            response.raise_for_status()
            file_bytes = BytesIO(response.content)
            return Presentation(file_bytes)
        return Presentation(self.document_path)

    def _get_shape_tag(self, shape) -> Optional[str]:
        """
        Extrae un tag como {TABLE_DATA} del campo 'name' de una forma (usado para tablas y gráficos).
        En PowerPoint, esto se establece en: Formato de Forma -> Panel de Selección -> Nombre.
        """
        try:
            descr = shape._element._nvXxPr.cNvPr.attrib.get("name", "")
            matches = re.findall(self.pattern, descr)
            if matches:
                return matches[0]
        except Exception as e:
            self.logger.error(f"Error al obtener tag de la forma: {e}")
        return None

    def analize_placeholders(self) -> Dict[str, Any]:
        """
        Analiza la presentación para encontrar todos los placeholders y los clasifica
        por tipo (string, table, chart).
        """
        prs = self.get_presentation()
        placeholders_encontrados = {}

        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            matches = re.findall(self.pattern, run.text)
                            for match in matches:
                                placeholders_encontrados[match] = {"type": "string"}

                if shape.has_table:
                    tag = self._get_shape_tag(shape)
                    if tag:
                        placeholders_encontrados[tag] = {
                            "type": "table",
                            "metadata": {
                                "row_number": len(shape.table.rows),
                                "column_number": len(shape.table.columns)
                            }
                        }

                if shape.has_chart:
                    tag = self._get_shape_tag(shape)
                    if tag:
                        placeholders_encontrados[tag] = {"type": "chart"}
                        
        return placeholders_encontrados

    def _update_table(self, shape, data: List[List[Any]]):

        table = shape.table
        
        for i, row_data in enumerate(data):
            if i < len(table.rows):
                for j, cell_value in enumerate(row_data):
                    if j < len(table.columns):
                        table.cell(i, j).text = str(cell_value)

    def _update_chart(self, shape, chart_data_json: Dict[str, Any]):
        chart = shape.chart
        chart_data = CategoryChartData()
        
        chart_data.categories = chart_data_json.get("categories", [])

        for serie in chart_data_json.get("series", []):
            numeric_values = [float(v) for v in serie.get("values", [])]
            chart_data.add_series(serie.get("name", ""), numeric_values)

        chart.replace_data(chart_data)

    def refill_document(self, fill_data: Dict[str, Any], output_path: Optional[str] = 'tmp/output.pptx') -> None:
        prs = self.get_presentation()

        for slide in prs.slides:
            for shape in slide.shapes:

                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for key, valor in fill_data.items():
                            placeholder = f"{{{{{key}}}}}"
                            if isinstance(valor, (str, int, float)):
                                if placeholder in paragraph.text:
                                    inline_text = paragraph.text
                                    paragraph.text = inline_text.replace(placeholder, str(valor))

                elif shape.has_table:
                    tag = self._get_shape_tag(shape)
                    if tag and tag in fill_data:
                        self._update_table(shape, fill_data[tag])
                
                elif shape.has_chart:
                    tag = self._get_shape_tag(shape)
                    if tag and tag in fill_data:
                        self._update_chart(shape, fill_data[tag])

        self.save_document(prs, output_path)

    def create_directories(self, file_path: str) -> None:
        from pathlib import Path
        path_object = Path(file_path)
        path_object.parent.mkdir(parents=True, exist_ok=True) 
        
    def save_document(self, document_content: PresentationObject, output_file_path: str) -> None:
        self.create_directories(output_file_path)
        document_content.save(output_file_path)

        self.logger.info(f"Documento guardado exitosamente en: {output_file_path}")

if __name__== '__main__':
    ppt_manager = PptManager('utility_test/Prueba.pptx')
    fill_data = ppt_manager.analize_placeholders()
    print("fill_data", fill_data)
    fill_data = {'TITULO_PRINCIPAL': 'Iniciativa de Data Analytics', 'NOMBRE_CLIENTE': 'Julio Romero', 'RESUMEN_EJECUTIVO': 'Estamos en un mundo lleno de dolor que se puede solucionar con data analytics', 'AUTOR': 'Julio Romero', 'INFORMATION_DATA': 'Datos relevantes sobre la problemática y su solución', 'TABLE_SCHEMA': [['Categoría', 'Datos Relevantes', 'Soluciones Propuestas'], ['Salud', 'Dolor Crónico', 'Análisis Predictivo'], ['Educación', 'Desigualdad', 'Programas de Intervención']], 'CHART_TITLE': 'Análisis de Impacto', 'CHART_SCHEMA': {'categories': ['Q1', 'Q2', 'Q3', 'Q4'], 'series': [{'name': 'Impacto 2024', 'values': [10, 20, 30, 40]}, {'name': 'Impacto 2025', 'values': [15, 25, 35, 45]}]}}
    ppt_manager.refill_document(fill_data)